from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta
from pptx.oxml import parse_xml
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree
import os

def set_pres_repeat(prs):
    prs_part = prs.part
    prs_props_part = prs_part.part_related_by(RT.PRES_PROPS)
    presentationPr = parse_xml(prs_props_part.blob)
    
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    
    # Find or create showPr element
    showPr_elements = presentationPr.findall(f'.//{{{p}}}showPr')
    if not showPr_elements:
        showPr_element = etree.Element(f'{{{p}}}showPr', loop="1", restart="always")
        presentationPr.append(showPr_element)
    else:
        for showPr in showPr_elements:
            showPr.set("loop", "1")
            showPr.set("restart", "always")
    
    prs_props_part._blob = etree.tostring(presentationPr)
    
    # Ensure slideShowPr is set to loop
    view_props_part = prs_part.part_related_by(RT.VIEW_PROPS)
    viewProps = parse_xml(view_props_part.blob)
    slideShowPr_elements = viewProps.findall(f'.//{{{p}}}slideShowPr')
    if not slideShowPr_elements:
        slideShowPr_element = etree.Element(f'{{{p}}}slideShowPr', loop="1")
        viewProps.append(slideShowPr_element)
    else:
        for slideShowPr in slideShowPr_elements:
            slideShowPr.set("loop", "1")
    
    view_props_part._blob = etree.tostring(viewProps)

def hide_last_slide(prs):
    slides = prs.slides
    if len(slides) > 0:
        last_slide = slides[-1]
        last_slide_element = last_slide.element
        last_slide_element.set('hidden', '1')
        print("Last slide hidden successfully.")
    else:
        print("No slides to hide.")

def create_presentation():
    prs = Presentation()
    set_pres_repeat(prs)

    today = datetime.now()
    tomorrow = today + timedelta(days=1)

    screenshots = [
        f'images/{today.strftime("%d-%m-%Y")}.png',
        f'images/{tomorrow.strftime("%d-%m-%Y")}.png',
        f'images/{today.strftime("%d-%m-%Y")}-fire.jpg',
        'images/placeholder.png'
    ]

    transition_xml_template = '''
        <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" advance="after" advTm="1000">
            <p:wipe />
        </p:transition>
    '''

    for i, screenshot in enumerate(screenshots):
        if not os.path.exists(screenshot):
            print(f"File not found: {screenshot}. Skipping this screenshot.")
            continue

        print(screenshot)
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        img_path = screenshot
        try:
            img = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

            left = int((Inches(10) - img.width) / 2)
            top = int((Inches(7.5) - img.height) / 2)

            img.left = left
            img.top = top
        except Exception as e:
            print(f"Error adding picture {img_path}: {e}")
            continue

        transition_xml = transition_xml_template
        transition_fragment = parse_xml(transition_xml)

        slide.element.append(transition_fragment)

        print(f"Transition applied to slide {i+1}")

    hide_last_slide(prs)
    prs.save('presentation.pptx')

if __name__ == "__main__":
    create_presentation()
